#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
为web界面测试创建PPT文件
"""

import os
from pptx import Presentation
from pptx.util import Inches

def create_test_ppt_for_web():
    """创建一个用于web界面测试的PPT文件"""
    prs = Presentation()
    
    # 第一张幻灯片 - 项目概述
    slide1 = prs.slides.add_slide(prs.slide_layouts[1])  # 标题和内容布局
    slide1.shapes.title.text = "盛远生产可视化项目导入规划"
    content1 = slide1.placeholders[1]
    tf1 = content1.text_frame
    tf1.text = """项目背景：
• 提升生产线数字化水平
• 建立实时监控体系
• 优化生产流程管理
• 实现数据驱动决策"""
    
    # 第二张幻灯片 - 技术架构
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "系统技术架构"
    content2 = slide2.placeholders[1]
    tf2 = content2.text_frame
    tf2.text = """核心组件：
1. 数据采集层 - 传感器数据收集
2. 数据处理层 - 实时数据分析
3. 可视化层 - 图表和仪表板
4. 应用层 - 用户交互界面
5. 存储层 - 历史数据管理"""
    
    # 第三张幻灯片 - 实施计划
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "项目实施计划"
    content3 = slide3.placeholders[1]
    tf3 = content3.text_frame
    tf3.text = """时间安排：
第一阶段（1-2月）：需求调研与系统设计
第二阶段（3-4月）：核心功能开发
第三阶段（5月）：系统集成与测试
第四阶段（6月）：部署上线与培训

预期效果：
• 生产效率提升20%
• 质量问题减少30%
• 决策响应时间缩短50%"""
    
    # 第四张幻灯片 - 风险控制
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "风险控制措施"
    content4 = slide4.placeholders[1]
    tf4 = content4.text_frame
    tf4.text = """主要风险点：
• 数据安全风险 → 加密传输和访问控制
• 系统稳定性风险 → 冗余设计和备份机制
• 人员培训风险 → 分阶段培训和文档支持
• 技术兼容性风险 → 标准化接口设计

应急预案：
• 建立24小时技术支持团队
• 制定系统故障快速恢复流程
• 准备备用系统和数据恢复方案"""
    
    # 保存到uploads目录
    output_dir = "/opt/work/kb/uploads/documents"
    os.makedirs(output_dir, exist_ok=True)
    
    # 使用当前时间戳命名
    from datetime import datetime
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    test_file = f"{output_dir}/{timestamp}_盛远生产可视化项目导入规划_测试.pptx"
    
    prs.save(test_file)
    print(f"测试PPT文件已创建: {test_file}")
    return test_file

if __name__ == "__main__":
    create_test_ppt_for_web()