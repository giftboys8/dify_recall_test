"""PowerPoint文档解析器

支持PPT/PPTX转换为PDF，以便在现有的PDF查看器中显示。
"""

import os
import tempfile
from pathlib import Path
from typing import Optional, Dict, Any
from io import BytesIO

# 尝试导入PPT处理依赖
Presentation = None
canvas = None
A4 = None

try:
    from pptx import Presentation
    print("Successfully imported pptx")
except ImportError as e:
    print(f"Warning: pptx not available: {e}")

try:
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.lib.units import inch
    print("Successfully imported reportlab")
except ImportError as e:
    print(f"Warning: reportlab not available: {e}")

from ..utils.logger import get_logger


class PPTParser:
    """PowerPoint文档解析器
    
    将PPT/PPTX文件转换为PDF格式，以便在现有的PDF查看器中显示。
    """
    
    def __init__(self, temp_dir: Optional[str] = None):
        """初始化PPT解析器
        
        Args:
            temp_dir: 临时文件目录，默认使用系统临时目录
        """
        self.temp_dir = temp_dir or tempfile.gettempdir()
        self.temp_files = []  # 跟踪临时文件用于清理
        self.logger = get_logger(__name__)
        
        if not self.is_available():
            self.logger.warning("PPT处理依赖未安装，PPT转换功能不可用")
    
    def is_available(self) -> bool:
        """检查PPT处理功能是否可用"""
        return all([
            Presentation is not None,
            canvas is not None,
            A4 is not None
        ])
    
    def can_convert_to_pdf(self) -> bool:
        """检查是否可以转换为PDF格式"""
        return all([
            Presentation is not None,
            canvas is not None,
            A4 is not None
        ])
    
    def convert_ppt_to_pdf(self, ppt_path: str, output_path: Optional[str] = None) -> str:
        """将PPT/PPTX文件转换为PDF（或文本文件作为备选方案）
        
        Args:
            ppt_path: PPT文件路径
            output_path: 输出文件路径，如果为None则自动生成
            
        Returns:
            生成的文件路径
            
        Raises:
            ValueError: 当PPT处理功能不可用时
            FileNotFoundError: 当输入文件不存在时
            Exception: 转换过程中的其他错误
        """
        if not os.path.exists(ppt_path):
            raise FileNotFoundError(f"PPT文件不存在: {ppt_path}")
        
        # 生成输出路径
        if output_path is None:
            base_name = Path(ppt_path).stem
            if self.can_convert_to_pdf():
                output_path = os.path.join(self.temp_dir, f"{base_name}_converted.pdf")
            else:
                output_path = os.path.join(self.temp_dir, f"{base_name}_converted.txt")
        
        try:
            self.logger.info(f"开始转换PPT: {ppt_path} -> {output_path}")
            
            if self.can_convert_to_pdf():
                # 使用完整的PDF转换功能
                return self._convert_to_pdf(ppt_path, output_path)
            else:
                # 使用文本提取备选方案
                return self._convert_to_text(ppt_path, output_path)
            
        except Exception as e:
            self.logger.error(f"PPT转换失败: {str(e)}")
            raise
    
    def _convert_to_pdf(self, ppt_path: str, output_path: str) -> str:
        """完整的PDF转换功能"""
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.units import inch
        
        # 加载PPT文件
        presentation = Presentation(ppt_path)
        
        # 创建PDF文档
        doc = SimpleDocTemplate(output_path, pagesize=A4)
        story = []
        styles = getSampleStyleSheet()
        
        # 尝试注册中文字体（如果可用）
        try:
            # 尝试使用系统中文字体
            import platform
            if platform.system() == "Darwin":  # macOS
                font_paths = [
                    "/System/Library/Fonts/PingFang.ttc",
                    "/System/Library/Fonts/Helvetica.ttc",
                    "/Library/Fonts/Arial Unicode MS.ttf"
                ]
            elif platform.system() == "Windows":
                font_paths = [
                    "C:/Windows/Fonts/msyh.ttc",  # 微软雅黑
                    "C:/Windows/Fonts/simsun.ttc",  # 宋体
                ]
            else:  # Linux
                font_paths = [
                    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
                    "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf"
                ]
            
            font_registered = False
            for font_path in font_paths:
                try:
                    if os.path.exists(font_path):
                        pdfmetrics.registerFont(TTFont('ChineseFont', font_path))
                        font_registered = True
                        break
                except:
                    continue
            
            if font_registered:
                styles['Normal'].fontName = 'ChineseFont'
                styles['Heading1'].fontName = 'ChineseFont'
        except Exception as e:
            self.logger.warning(f"无法注册中文字体，将使用默认字体: {e}")
        
        # 处理每张幻灯片
        for slide_idx, slide in enumerate(presentation.slides):
            self.logger.debug(f"处理幻灯片 {slide_idx + 1}/{len(presentation.slides)}")
            
            # 添加幻灯片标题
            title = f"幻灯片 {slide_idx + 1}"
            story.append(Paragraph(title, styles['Heading1']))
            story.append(Spacer(1, 0.2*inch))
            
            # 提取并添加文本内容
            slide_content = []
            for shape in slide.shapes:
                try:
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text.strip()
                        # 处理文本格式
                        text = text.replace('\n', '<br/>')
                        slide_content.append(text)
                except Exception as e:
                    # 跳过无法识别的形状类型
                    self.logger.warning(f"跳过无法处理的形状: {e}")
                    continue
            
            if slide_content:
                for content in slide_content:
                    try:
                        para = Paragraph(content, styles['Normal'])
                        story.append(para)
                        story.append(Spacer(1, 0.1*inch))
                    except Exception as e:
                        # 如果Paragraph创建失败，使用纯文本
                        self.logger.warning(f"文本处理失败，使用纯文本: {e}")
                        story.append(Paragraph(content.replace('<br/>', ' '), styles['Normal']))
                        story.append(Spacer(1, 0.1*inch))
            else:
                story.append(Paragraph("(此幻灯片无文本内容)", styles['Normal']))
                story.append(Spacer(1, 0.1*inch))
            
            # 添加图片信息提示
            image_count = 0
            for shape in slide.shapes:
                try:
                    if hasattr(shape, 'shape_type') and 'PICTURE' in str(shape.shape_type):
                        image_count += 1
                except Exception as e:
                    # 跳过无法识别的形状类型
                    self.logger.debug(f"跳过无法识别的形状类型: {e}")
                    continue
            
            if image_count > 0:
                story.append(Paragraph(f"[此幻灯片包含 {image_count} 张图片，PDF转换中图片暂不支持]", styles['Normal']))
                story.append(Spacer(1, 0.1*inch))
            
            # 添加分隔线
            if slide_idx < len(presentation.slides) - 1:
                story.append(Spacer(1, 0.3*inch))
                story.append(Paragraph("_" * 50, styles['Normal']))
                story.append(Spacer(1, 0.3*inch))
        
        # 构建PDF
        try:
            doc.build(story)
        except Exception as e:
            self.logger.error(f"PDF构建失败: {e}")
            # 回退到简单的文本方式
            return self._convert_to_text(ppt_path, output_path.replace('.pdf', '.txt'))
        
        # 记录临时文件
        if output_path.startswith(self.temp_dir):
            self.temp_files.append(output_path)
        
        self.logger.info(f"PPT转换完成: {output_path}")
        return output_path
    
    def _convert_to_text(self, ppt_path: str, output_path: str) -> str:
        """文本提取备选方案"""
        if Presentation is None:
            # 如果连pptx都没有，创建一个简单的错误文件
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write("PPT处理功能不可用\n")
                f.write("请安装必要的依赖包：python-pptx, reportlab, Pillow\n")
                f.write(f"原文件: {ppt_path}\n")
            return output_path
        
        # 加载PPT文件
        presentation = Presentation(ppt_path)
        
        # 提取文本内容
        content = []
        content.append(f"PPT文档内容提取\n")
        content.append(f"原文件: {ppt_path}\n")
        content.append(f"幻灯片总数: {len(presentation.slides)}\n")
        content.append("=" * 50 + "\n")
        
        for slide_idx, slide in enumerate(presentation.slides):
            content.append(f"\n幻灯片 {slide_idx + 1}:\n")
            content.append("-" * 30 + "\n")
            
            # 提取文本内容
            slide_text = []
            for shape in slide.shapes:
                try:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                except Exception as e:
                    # 跳过无法识别的形状类型
                    self.logger.warning(f"跳过无法处理的形状: {e}")
                    continue
            
            if slide_text:
                content.extend([text + "\n" for text in slide_text])
            else:
                content.append("(此幻灯片无文本内容)\n")
        
        # 保存文本文件
        with open(output_path, 'w', encoding='utf-8') as f:
            f.writelines(content)
        
        # 记录临时文件
        if output_path.startswith(self.temp_dir):
            self.temp_files.append(output_path)
        
        self.logger.info(f"PPT文本提取完成: {output_path}")
        return output_path
    
    def extract_slide_info(self, ppt_path: str) -> Dict[str, Any]:
        """提取PPT文件的基本信息
        
        Args:
            ppt_path: PPT文件路径
            
        Returns:
            包含PPT信息的字典
        """
        if not self.is_available():
            return {"error": "PPT处理功能不可用"}
        
        try:
            presentation = Presentation(ppt_path)
            
            info = {
                "slide_count": len(presentation.slides),
                "slide_size": {
                    "width": presentation.slide_width,
                    "height": presentation.slide_height
                },
                "slides": []
            }
            
            for idx, slide in enumerate(presentation.slides):
                slide_info = {
                    "index": idx + 1,
                    "text_content": [],
                    "shape_count": len(slide.shapes)
                }
                
                # 提取文本内容
                for shape in slide.shapes:
                    try:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_info["text_content"].append(shape.text.strip())
                    except Exception as e:
                        # 跳过无法识别的形状类型
                        self.logger.debug(f"跳过无法识别的形状类型: {e}")
                        continue
                
                info["slides"].append(slide_info)
            
            return info
            
        except Exception as e:
            self.logger.error(f"提取PPT信息失败: {str(e)}")
            return {"error": str(e)}
    
    def cleanup(self):
        """清理临时文件"""
        for temp_file in self.temp_files:
            try:
                if os.path.exists(temp_file):
                    os.remove(temp_file)
                    self.logger.debug(f"已删除临时文件: {temp_file}")
            except Exception as e:
                self.logger.warning(f"删除临时文件失败 {temp_file}: {str(e)}")
        
        self.temp_files.clear()
    
    def __del__(self):
        """析构函数，自动清理临时文件"""
        self.cleanup()