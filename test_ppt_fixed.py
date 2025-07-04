#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
测试修复后的PPT处理功能
"""

import os
import sys
import tempfile
from pptx import Presentation
from pptx.util import Inches

# 添加项目路径
sys.path.append('/opt/work/kb')
from src.translation.ppt_parser import PPTParser

def create_test_ppt():
    """创建一个包含中文内容的测试PPT文件"""
    prs = Presentation()
    
    # 第一张幻灯片
    slide1 = prs.slides.add_slide(prs.slide_layouts[1])  # 标题和内容布局
    slide1.shapes.title.text = "盛远生产可视化项目导入规划"
    content1 = slide1.placeholders[1]
    tf1 = content1.text_frame
    tf1.text = "项目概述：\n• 建立生产数据可视化平台\n• 提升生产效率和质量管控\n• 实现智能化生产管理"
    
    # 第二张幻灯片
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "技术架构"
    content2 = slide2.placeholders[1]
    tf2 = content2.text_frame
    tf2.text = "系统组成：\n1. 数据采集层\n2. 数据处理层\n3. 可视化展示层\n4. 用户交互层"
    
    # 第三张幻灯片
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "实施计划"
    content3 = slide3.placeholders[1]
    tf3 = content3.text_frame
    tf3.text = "阶段安排：\n第一阶段：需求分析和系统设计\n第二阶段：核心功能开发\n第三阶段：测试和部署\n第四阶段：培训和上线"
    
    # 保存文件
    test_file = "/tmp/test_chinese_ppt.pptx"
    prs.save(test_file)
    print(f"测试PPT文件已创建: {test_file}")
    return test_file

def test_ppt_conversion():
    """测试PPT转换功能"""
    print("=== 测试修复后的PPT处理功能 ===")
    
    # 创建测试PPT
    test_ppt = create_test_ppt()
    
    try:
        # 初始化PPT解析器
        parser = PPTParser()
        print(f"PPT解析器初始化成功")
        
        # 测试可用性
        if parser.is_available():
            print("✓ PPT解析器可用")
        else:
            print("✗ PPT解析器不可用")
            return
        
        # 测试PDF转换
        if parser.can_convert_to_pdf():
            print("✓ 支持PDF转换")
        else:
            print("✗ 不支持PDF转换")
            return
        
        # 执行转换
        print("\n开始转换PPT到PDF...")
        output_file = parser.convert_ppt_to_pdf(test_ppt)
        
        if output_file and os.path.exists(output_file):
            print(f"✓ 转换成功: {output_file}")
            
            # 检查文件大小
            file_size = os.path.getsize(output_file)
            print(f"输出文件大小: {file_size} 字节")
            
            if output_file.endswith('.pdf'):
                print("✓ 成功生成PDF文件")
                
                # 尝试读取PDF内容进行验证
                try:
                    from PyPDF2 import PdfReader
                    reader = PdfReader(output_file)
                    total_text = ""
                    for page in reader.pages:
                        total_text += page.extract_text()
                    
                    print(f"PDF页数: {len(reader.pages)}")
                    print(f"提取的文本长度: {len(total_text)} 字符")
                    
                    # 检查是否包含中文
                    chinese_chars = sum(1 for char in total_text if '\u4e00' <= char <= '\u9fff')
                    print(f"中文字符数量: {chinese_chars}")
                    
                    if chinese_chars > 0:
                        print("✓ 中文字符处理正常")
                    else:
                        print("⚠ 未检测到中文字符")
                    
                    # 显示部分内容
                    if total_text:
                        content_preview = total_text[:300].replace('\n', ' ')
                        print(f"\n内容预览: {content_preview}...")
                        
                except ImportError:
                    print("PyPDF2未安装，无法验证PDF内容")
                except Exception as e:
                    print(f"PDF内容验证失败: {e}")
            else:
                print(f"生成了文本文件: {output_file}")
                
            # 清理输出文件
            os.remove(output_file)
            print(f"输出文件已清理: {output_file}")
                
        else:
            print("✗ 转换失败")
            
    except Exception as e:
        print(f"✗ 测试过程中出现错误: {e}")
        import traceback
        traceback.print_exc()
    
    finally:
        # 清理测试文件
        if os.path.exists(test_ppt):
            os.remove(test_ppt)
            print(f"\n测试文件已清理: {test_ppt}")

if __name__ == "__main__":
    test_ppt_conversion()